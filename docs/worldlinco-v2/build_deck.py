# -*- coding: utf-8 -*-
"""WorldLinco 사업용 PowerPoint 덱 생성기.

소리새 하늘색 배경 룩으로 표지 / 기능 / 화면 / 사용법 / 구조도 / 사업 설명서를 묶는다.
실행: python docs/worldlinco-v2/build_deck.py
출력: docs/worldlinco-v2/WorldLinco_사업소개.pptx
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"
OUT = HERE / "WorldLinco_사업소개.pptx"

# --- 디자인 토큰 (theme.ts 정합) ---
SKY_TOP = RGBColor(0xE3, 0xF0, 0xFF)
SKY_MID = RGBColor(0xF0, 0xF7, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PRIMARY = RGBColor(0x1E, 0x6F, 0xE0)
PRIMARY_DK = RGBColor(0x0B, 0x2E, 0x5E)
CORAL = RGBColor(0xFF, 0x8A, 0x5B)
SUCCESS = RGBColor(0x19, 0xC3, 0x7D)
PURPLE = RGBColor(0x7C, 0x5C, 0xFC)
TEXT = RGBColor(0x1A, 0x1F, 0x36)
MUTED = RGBColor(0x6B, 0x72, 0x80)
CARD_LINE = RGBColor(0xDF, 0xE7, 0xF2)

KR_FONT = "맑은 고딕"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)


def _set_kr(run, *, size, bold=False, color=TEXT):
    run.font.name = KR_FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # 한글 글꼴 강제(eastasia)
    rPr = run._r.get_or_add_rPr()
    import copy
    latin = rPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}latin')
    ea = rPr.makeelement('{http://schemas.openxmlformats.org/drawingml/2006/main}ea', {'typeface': KR_FONT})
    rPr.append(ea)


def sky_background(slide):
    """슬라이드 배경을 소리새 하늘색 세로 그라데이션으로."""
    fill = slide.background.fill
    fill.gradient()
    try:
        stops = fill.gradient_stops
        stops[0].position = 0.0
        stops[0].color.rgb = SKY_TOP
        stops[1].position = 1.0
        stops[1].color.rgb = WHITE
        # 중간 흰빛
        fill.gradient_angle = 90.0
    except Exception:
        fill.solid()
        fill.fore_color.rgb = SKY_MID


def add_text(slide, left, top, width, height, lines, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, bold, color) tuples → 단락별."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        _set_kr(run, size=size, bold=bold, color=color)
    return tb


def add_card(slide, left, top, width, height, *, fill=WHITE, line=CARD_LINE,
             radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    box = slide.shapes.add_shape(shape_type, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(1)
    box.shadow.inherit = False
    return box


def add_image_fit(slide, path: Path, left, top, width, height):
    """박스 안에 비율 유지하며 가운데 정렬."""
    iw, ih = Image.open(path).size
    box_ratio = width / height
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        w = width
        h = int(width / img_ratio)
    else:
        h = height
        w = int(height * img_ratio)
    x = left + (width - w) // 2
    y = top + (height - h) // 2
    return slide.shapes.add_picture(str(path), x, y, w, h)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    sky_background(slide)
    return slide


def chip(slide, left, top, text, color):
    c = add_card(slide, left, top, Inches(2.2), Inches(0.42), fill=WHITE, line=color)
    tf = c.text_frame
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    _set_kr(r, size=12, bold=True, color=color)
    return c


def main():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H

    # ---------- 1. 표지 ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # 풀블리드 광고 표지
    add_image_fit(s, ASSETS / "worldlinco_ad_cover.png", 0, 0, EMU_W, EMU_H)

    # ---------- 2. 한눈에 ----------
    s = new_slide(prs)
    add_text(s, Inches(0.8), Inches(0.7), Inches(11.7), Inches(1.0), [
        ("WorldLinco — 언어의 벽을 넘어, 세상을 연결합니다", 30, True, PRIMARY_DK),
    ])
    add_text(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.6), [
        ("실시간 AI 통역 슈퍼앱 · 대면·통화·채팅·노래·여행을 하나로", 16, False, MUTED),
    ])
    pain = [
        ("문제", 18, True, CORAL),
        ("· 해외 여행·외국인 응대 시 언어 장벽으로 소통 단절", 15, False, TEXT),
        ("· 통역 앱은 분절적(번역기 ≠ 통화 ≠ 채팅)이고 지연이 큼", 15, False, TEXT),
        ("· 개인화된 통역 비서 부재", 15, False, TEXT),
    ]
    sol = [
        ("해결 — WorldLinco", 18, True, SUCCESS),
        ("· 초저지연 실시간 통역 파이프라인", 15, False, TEXT),
        ("· 대면·VoIP·채팅·노래·예약 통합 채널", 15, False, TEXT),
        ("· 소리새: 기억·성격을 가진 개인 AI 동반자", 15, False, TEXT),
    ]
    add_card(s, Inches(0.8), Inches(2.5), Inches(5.6), Inches(3.6))
    add_text(s, Inches(1.1), Inches(2.8), Inches(5.0), Inches(3.0), pain)
    add_card(s, Inches(6.9), Inches(2.5), Inches(5.6), Inches(3.6))
    add_text(s, Inches(7.2), Inches(2.8), Inches(5.0), Inches(3.0), sol)

    # ---------- 3. 핵심 기능 ----------
    s = new_slide(prs)
    add_text(s, Inches(0.8), Inches(0.55), Inches(11.7), Inches(0.9), [
        ("핵심 기능", 30, True, PRIMARY_DK),
    ])
    features = [
        ("① 대면 통역", "마주보고 자동 언어감지 · 화면 분할 양방향 실시간 통역", PRIMARY),
        ("② 소리새 AI 동반자", "웨이크워드 호출 · 기억·성격을 가진 개인 AI 통역 비서", PRIMARY),
        ("③ VoIP 통역 통화", "인터넷 음성통화에 실시간 자막 통역(원문+번역)", PRIMARY_DK),
        ("④ 채팅 번역", "텍스트·음성 겸용 · 번역 말풍선 페어", PRIMARY),
        ("⑤ 노래 번역", "가사 라인별 원문/번역 동기화 재생", PURPLE),
        ("⑥ 예약·일반전화", "항공·호텔 예약 · PSTN 통화 보조", SUCCESS),
    ]
    gx, gy = Inches(0.8), Inches(1.6)
    cw, ch = Inches(3.78), Inches(2.45)
    gap = Inches(0.2)
    for idx, (title, desc, color) in enumerate(features):
        r, c = divmod(idx, 3)
        left = Emu(int(gx) + c * (int(cw) + int(gap)))
        top = Emu(int(gy) + r * (int(ch) + int(gap)))
        add_card(s, left, top, cw, ch)
        bar = add_card(s, left, top, Inches(0.14), ch, fill=color, line=color)
        add_text(s, Emu(int(left) + Inches(0.4)), Emu(int(top) + Inches(0.3)),
                 Emu(int(cw) - Inches(0.6)), Inches(0.6), [(title, 17, True, color)])
        add_text(s, Emu(int(left) + Inches(0.4)), Emu(int(top) + Inches(1.0)),
                 Emu(int(cw) - Inches(0.6)), Inches(1.3), [(desc, 13, False, TEXT)])

    # ---------- 4. 통합 화면 시안 ----------
    s = new_slide(prs)
    add_text(s, Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.8), [
        ("하나의 디자인 — 소리새 하늘색 통합 UI", 26, True, PRIMARY_DK),
    ])
    add_image_fit(s, ASSETS / "worldlinco_unified_sky_bg.png",
                  Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.7))

    # ---------- 5. 주요 화면 (단일 목업) ----------
    s = new_slide(prs)
    add_text(s, Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.8), [
        ("주요 화면", 30, True, PRIMARY_DK),
    ])
    shots = [
        ("worldlinco_screen_voip_sky.png", "VoIP 통역 통화"),
        ("worldlinco_screen_song_sky.png", "노래 번역"),
        ("worldlinco_screen_booking_sky.png", "여행 예약"),
    ]
    sw = Inches(4.0)
    sx = Inches(0.45)
    for i, (fn, cap) in enumerate(shots):
        left = Emu(int(sx) + i * (int(sw) + int(Inches(0.2))))
        add_image_fit(s, ASSETS / fn, left, Inches(1.4), sw, Inches(4.7))
        add_text(s, left, Inches(6.2), sw, Inches(0.5),
                 [(cap, 15, True, PRIMARY_DK)], align=PP_ALIGN.CENTER)

    # ---------- 6. 사용 설명서 ----------
    s = new_slide(prs)
    add_text(s, Inches(0.8), Inches(0.55), Inches(11.7), Inches(0.9), [
        ("사용 설명서 — 3단계로 통역 시작", 30, True, PRIMARY_DK),
    ])
    steps = [
        ("1. 언어 선택", "내 언어와 상대 언어를 선택하거나 대면통역에서 자동 감지를 켭니다."),
        ("2. 말하기 / 입력", "원형 마이크 버튼을 누르고 말하거나 채팅에 입력합니다."),
        ("3. 실시간 통역 확인", "원문과 번역이 쌍으로 표시되고 음성으로 읽어줍니다."),
    ]
    yy = Inches(1.8)
    for i, (t, d) in enumerate(steps):
        top = Emu(int(yy) + i * int(Inches(1.55)))
        num = add_card(s, Inches(0.9), top, Inches(0.9), Inches(0.9),
                       fill=PRIMARY, line=PRIMARY)
        ntf = num.text_frame
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        np = ntf.paragraphs[0]
        np.alignment = PP_ALIGN.CENTER
        nr = np.add_run()
        nr.text = str(i + 1)
        _set_kr(nr, size=26, bold=True, color=WHITE)
        add_text(s, Inches(2.1), Emu(int(top) - int(Inches(0.05))), Inches(10), Inches(0.6),
                 [(t, 19, True, PRIMARY_DK)])
        add_text(s, Inches(2.1), Emu(int(top) + int(Inches(0.5))), Inches(10), Inches(0.8),
                 [(d, 14, False, TEXT)])

    # ---------- 7. 구조도 ----------
    s = new_slide(prs)
    add_text(s, Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.8), [
        ("시스템 구조도", 30, True, PRIMARY_DK),
    ])
    layers = [
        ("클라이언트 (모바일 앱 · React Native / Expo)", PRIMARY),
        ("API 게이트웨이 (인증 · 라우팅 · 레이트리밋)", PRIMARY),
        ("통신 오케스트레이터 (세션 · 채널 정책)", PRIMARY_DK),
        ("기능 허브 — 대면 · VoIP · 채팅 · 노래 · 예약 · 소리새", PRIMARY_DK),
        ("지능 엔진 (STT · NMT · TTS · LLM 파이프라인)", PURPLE),
        ("음성 파이프라인 (저지연 스트리밍 · WebRTC)", PURPLE),
        ("인프라 (GPU 컴퓨트 · 저장소 · 메시징 · 모니터링)", PRIMARY_DK),
    ]
    ly = Inches(1.5)
    lh = Inches(0.72)
    lgap = Inches(0.12)
    for i, (label, color) in enumerate(layers):
        top = Emu(int(ly) + i * (int(lh) + int(lgap)))
        box = add_card(s, Inches(1.8), top, Inches(9.7), lh, fill=WHITE, line=color)
        bar = add_card(s, Inches(1.8), top, Inches(0.16), lh, fill=color, line=color)
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        rr = p.add_run()
        rr.text = label
        _set_kr(rr, size=15, bold=True, color=TEXT)
        if i < len(layers) - 1:
            arrow = s.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW, Inches(6.55),
                Emu(int(top) + int(lh) - int(Inches(0.02))),
                Inches(0.25), Inches(0.14))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MUTED
            arrow.line.fill.background()

    # ---------- 8. 사업 모델 ----------
    s = new_slide(prs)
    add_text(s, Inches(0.8), Inches(0.55), Inches(11.7), Inches(0.9), [
        ("사업 설명 — 시장 · 수익 모델", 30, True, PRIMARY_DK),
    ])
    add_card(s, Inches(0.8), Inches(1.7), Inches(5.6), Inches(4.6))
    add_text(s, Inches(1.1), Inches(1.95), Inches(5.0), Inches(4.2), [
        ("시장 기회", 19, True, PRIMARY),
        ("· 글로벌 여행·관광 통역 수요 급증", 14, False, TEXT),
        ("· 외국인 거주·의료·행정 응대 통역", 14, False, TEXT),
        ("· 비대면 글로벌 커뮤니케이션 확대", 14, False, TEXT),
        ("", 8, False, TEXT),
        ("차별화 포인트", 19, True, CORAL),
        ("· 초저지연 실시간 통역 파이프라인", 14, False, TEXT),
        ("· 다채널(대면/통화/채팅) 단일 앱", 14, False, TEXT),
        ("· 소리새 = 개인 AI 통역 동반자", 14, False, TEXT),
    ])
    add_card(s, Inches(6.9), Inches(1.7), Inches(5.6), Inches(4.6))
    add_text(s, Inches(7.2), Inches(1.95), Inches(5.0), Inches(4.2), [
        ("수익 모델", 19, True, SUCCESS),
        ("· Free — 기본 번역·일 사용량 제한", 14, False, TEXT),
        ("· Pro(개인 구독) — 무제한·고품질 음성", 14, False, TEXT),
        ("· Business — 기업·여행사·콜센터", 14, False, TEXT),
        ("· API/B2B — 통역 엔진 라이선스", 14, False, TEXT),
        ("", 8, False, TEXT),
        ("로드맵", 19, True, PRIMARY),
        ("· 통역 품질·언어 확장 → 소리새 고도화", 14, False, TEXT),
        ("· B2B 채널(여행/의료/공공) 확대", 14, False, TEXT),
    ])

    # ---------- 9. 마무리 ----------
    s = new_slide(prs)
    add_text(s, Inches(0), Inches(2.6), EMU_W, Inches(1.2), [
        ("WorldLinco", 48, True, PRIMARY_DK),
    ], align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(3.9), EMU_W, Inches(0.8), [
        ("언어의 벽을 넘어, 세상을 연결합니다", 22, False, PRIMARY),
    ], align=PP_ALIGN.CENTER)

    prs.save(OUT)
    print(f"saved: {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
